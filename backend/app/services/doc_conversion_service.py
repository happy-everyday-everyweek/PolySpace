from __future__ import annotations

import logging
import os
import tempfile
import uuid
from typing import Any

from app.services.libreoffice_service import libreoffice_service

logger = logging.getLogger(__name__)

UPLOAD_DIR = os.path.join(os.getcwd(), "data", "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)


def _resolve_path(file_path: str) -> str:
    if os.path.isabs(file_path) and os.path.isfile(file_path):
        return file_path
    candidate = os.path.join(UPLOAD_DIR, file_path)
    if os.path.isfile(candidate):
        return candidate
    return file_path


def _output_path(suffix: str) -> str:
    file_id = str(uuid.uuid4()) + suffix
    return os.path.join(UPLOAD_DIR, file_id), file_id


class DocumentConversionService:

    async def convert_file(
        self,
        input_path: str,
        output_format: str,
        source_format: str | None = None,
    ) -> dict[str, Any]:
        resolved = _resolve_path(input_path)
        if not os.path.isfile(resolved):
            return {"error": f"File not found: {input_path}"}

        if source_format is None:
            ext = os.path.splitext(resolved)[1].lower().lstrip(".")
            source_format = ext or "unknown"

        if output_format == source_format:
            return {"error": f"Source and target format are the same: {output_format}"}

        if source_format == "html" and output_format in ("docx", "odt", "pdf", "rtf", "txt"):
            return await self._html_to_office(resolved, output_format)

        if source_format in ("docx", "odt", "rtf", "doc") and output_format in ("pdf", "html", "txt", "odt", "docx"):
            return await self._office_convert(resolved, output_format)

        if source_format in ("pptx", "odp", "ppt") and output_format in ("pdf", "odp", "pptx"):
            return await self._office_convert(resolved, output_format)

        if source_format in ("xlsx", "ods", "xls", "csv") and output_format in ("pdf", "ods", "xlsx", "csv", "html"):
            return await self._office_convert(resolved, output_format)

        if source_format == "pdf" and output_format in ("docx", "html", "txt", "odt"):
            return await self._pdf_to_office(resolved, output_format)

        return await self._office_convert(resolved, output_format)

    async def html_to_docx(
        self,
        html_content: str,
        title: str = "Document",
    ) -> dict[str, Any]:
        with tempfile.TemporaryDirectory(dir=UPLOAD_DIR) as tmpdir:
            html_path = os.path.join(tmpdir, "input.html")
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(self._wrap_html(html_content, title))
            return await self._html_to_office(html_path, "docx")

    async def html_to_pdf(
        self,
        html_content: str,
        title: str = "Document",
    ) -> dict[str, Any]:
        with tempfile.TemporaryDirectory(dir=UPLOAD_DIR) as tmpdir:
            html_path = os.path.join(tmpdir, "input.html")
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(self._wrap_html(html_content, title))
            return await self._html_to_office(html_path, "pdf")

    async def spreadsheet_data_to_xlsx(
        self,
        data: list[list[Any]],
        title: str = "Spreadsheet",
        headers: list[str] | None = None,
    ) -> dict[str, Any]:
        try:
            from openpyxl import Workbook

            wb = Workbook()
            ws = wb.active
            ws.title = title[:31]

            if headers:
                ws.append(headers)
            for row in data:
                ws.append(row)

            output_path, file_id = _output_path(".xlsx")
            wb.save(output_path)
            return {
                "status": "ok",
                "output_file_id": file_id,
                "output_path": output_path,
                "output_format": "xlsx",
                "file_size": os.path.getsize(output_path),
            }
        except ImportError:
            return {"error": "openpyxl not installed. Run: pip install openpyxl"}
        except Exception as e:
            return {"error": str(e)}

    async def spreadsheet_data_to_csv(
        self,
        data: list[list[Any]],
        headers: list[str] | None = None,
    ) -> dict[str, Any]:
        import csv
        output_path, file_id = _output_path(".csv")
        with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
            writer = csv.writer(f)
            if headers:
                writer.writerow(headers)
            for row in data:
                writer.writerow(row)
        return {
            "status": "ok",
            "output_file_id": file_id,
            "output_path": output_path,
            "output_format": "csv",
            "file_size": os.path.getsize(output_path),
        }

    async def slides_data_to_pptx(
        self,
        slides: list[dict[str, Any]],
        title: str = "Presentation",
    ) -> dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for slide_data in slides:
                layout_name = slide_data.get("layout", "title_content")
                if layout_name == "title":
                    slide_layout = prs.slide_layouts[0]
                elif layout_name == "blank":
                    slide_layout = prs.slide_layouts[6]
                else:
                    slide_layout = prs.slide_layouts[1]

                slide = prs.slides.add_slide(slide_layout)

                if slide_layout.has_placeholder(0):
                    title_placeholder = slide.placeholders[0]
                    title_placeholder.text = slide_data.get("title", "")

                if layout_name != "title" and slide_layout.has_placeholder(1):
                    body_placeholder = slide.placeholders[1]
                    bullets = slide_data.get("bullets", [])
                    if bullets:
                        tf = body_placeholder.text_frame
                        tf.clear()
                        for i, bullet in enumerate(bullets):
                            if isinstance(bullet, str):
                                if i == 0:
                                    tf.paragraphs[0].text = bullet
                                else:
                                    tf.add_paragraph().text = bullet
                            elif isinstance(bullet, dict):
                                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                                p.text = bullet.get("text", "")

                notes = slide_data.get("notes", "")
                if notes:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes

            output_path, file_id = _output_path(".pptx")
            prs.save(output_path)
            return {
                "status": "ok",
                "output_file_id": file_id,
                "output_path": output_path,
                "output_format": "pptx",
                "file_size": os.path.getsize(output_path),
            }
        except ImportError:
            return {"error": "python-pptx not installed. Run: pip install python-pptx"}
        except Exception as e:
            logger.error(f"PPTX generation error: {e}")
            return {"error": str(e)}

    async def _html_to_office(self, html_path: str, output_format: str) -> dict[str, Any]:
        if libreoffice_service.available:
            output_path, file_id = _output_path(f".{output_format}")
            result = await libreoffice_service.convert(html_path, output_format, output_path)
            if "error" not in result:
                result["output_file_id"] = file_id
                return result
            logger.warning(f"LibreOffice conversion failed, trying fallback: {result['error']}")

        if output_format == "docx":
            return await self._html_to_docx_fallback(html_path)
        if output_format == "pdf":
            return {"error": "PDF conversion requires LibreOffice. Please install LibreOffice."}

        return {"error": f"Conversion from HTML to {output_format} requires LibreOffice"}

    async def _html_to_docx_fallback(self, html_path: str) -> dict[str, Any]:
        try:
            from bs4 import BeautifulSoup
            from docx import Document

            with open(html_path, "r", encoding="utf-8") as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, "html.parser")
            doc = Document()

            for h in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                level = int(h.name[1])
                text = h.get_text(strip=True)
                if text:
                    doc.add_heading(text, level=level)

            for p in soup.find_all("p"):
                text = p.get_text(strip=True)
                if text:
                    doc.add_paragraph(text)

            for li in soup.find_all("li"):
                text = li.get_text(strip=True)
                if text:
                    doc.add_paragraph(text, style="List Bullet")

            for table in soup.find_all("table"):
                rows = table.find_all("tr")
                if not rows:
                    continue
                first_row = rows[0]
                cols = first_row.find_all(["td", "th"])
                table_obj = doc.add_table(rows=len(rows), cols=len(cols))
                table_obj.style = "Table Grid"
                for i, row in enumerate(rows):
                    cells = row.find_all(["td", "th"])
                    for j, cell in enumerate(cells):
                        if j < len(cols):
                            table_obj.rows[i].cells[j].text = cell.get_text(strip=True)

            output_path, file_id = _output_path(".docx")
            doc.save(output_path)
            return {
                "status": "ok",
                "output_file_id": file_id,
                "output_path": output_path,
                "output_format": "docx",
                "file_size": os.path.getsize(output_path),
                "fallback": True,
            }
        except ImportError as e:
            return {"error": f"Required library not installed: {e}"}
        except Exception as e:
            return {"error": str(e)}

    async def _office_convert(self, input_path: str, output_format: str) -> dict[str, Any]:
        if not libreoffice_service.available:
            return {"error": f"Conversion requires LibreOffice. Please install LibreOffice for {output_format} output."}

        output_path, file_id = _output_path(f".{output_format}")
        result = await libreoffice_service.convert(input_path, output_format, output_path)
        if "error" not in result:
            result["output_file_id"] = file_id
        return result

    async def _pdf_to_office(self, pdf_path: str, output_format: str) -> dict[str, Any]:
        if not libreoffice_service.available:
            if output_format == "txt":
                return await self._pdf_to_text_fallback(pdf_path)
            return {"error": f"PDF to {output_format} conversion requires LibreOffice"}

        output_path, file_id = _output_path(f".{output_format}")
        result = await libreoffice_service.convert(pdf_path, output_format, output_path)
        if "error" not in result:
            result["output_file_id"] = file_id
        return result

    async def _pdf_to_text_fallback(self, pdf_path: str) -> dict[str, Any]:
        try:
            import fitz
            doc = fitz.open(pdf_path)
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            doc.close()

            output_path, file_id = _output_path(".txt")
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n\n".join(text_parts))
            return {
                "status": "ok",
                "output_file_id": file_id,
                "output_path": output_path,
                "output_format": "txt",
                "file_size": os.path.getsize(output_path),
                "fallback": True,
            }
        except ImportError:
            return {"error": "PyMuPDF not installed for PDF text extraction"}
        except Exception as e:
            return {"error": str(e)}

    def _wrap_html(self, html_content: str, title: str = "Document") -> str:
        if "<html" in html_content.lower():
            return html_content
        return f"""<!DOCTYPE html>
<html>
<head><meta charset="utf-8"><title>{title}</title>
<style>
body {{ font-family: "Microsoft YaHei", "SimSun", Arial, sans-serif; margin: 2cm; }}
h1 {{ font-size: 24pt; font-weight: bold; }}
h2 {{ font-size: 18pt; font-weight: bold; }}
h3 {{ font-size: 14pt; font-weight: bold; }}
p {{ font-size: 12pt; line-height: 1.6; }}
table {{ border-collapse: collapse; width: 100%; }}
td, th {{ border: 1px solid #ccc; padding: 6px 10px; }}
th {{ background: #f0f0f0; font-weight: bold; }}
</style>
</head>
<body>
{html_content}
</body>
</html>"""


doc_conversion_service = DocumentConversionService()
